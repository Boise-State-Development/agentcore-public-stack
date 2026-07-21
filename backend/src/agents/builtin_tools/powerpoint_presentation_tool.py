"""PowerPoint presentation tools (create / modify / list / read).

Each tool runs python-pptx code inside AWS Bedrock Code Interpreter and uses the
existing user-files store (``apis.shared.files``) for persistence and delivery
— generated/modified ``.pptx`` files land in ``S3_USER_FILES_BUCKET_NAME`` with
a ``FileMetadata`` row (status READY) in ``DYNAMODB_USER_FILES_TABLE_NAME``, so
they appear in the chat's Files panel and are downloadable via the app-api
``/files/{id}/preview-url`` route.

Tools
-----
* ``create_powerpoint_presentation`` — build a new deck from python-pptx code.
* ``modify_powerpoint_presentation`` — edit an existing deck with python-pptx.
* ``list_powerpoint_presentations``  — list the .pptx files in this chat.
* ``read_powerpoint_presentation``   — extract a deck's slide text + notes.

(A slide-screenshot/preview tool is intentionally omitted: rasterizing a .pptx
requires LibreOffice/poppler, which the Python-only Code Interpreter sandbox
does not provide — the same reason the Word toolset omits one.)

Design notes
------------
* The Code Interpreter + user-files storage plumbing is shared with the Word
  and Excel toolsets and lives in ``builtin_tools.office._storage``; this module
  keeps only the python-pptx specifics (preamble, generate/modify/extract) and
  the four tool factories.
* Identity (``user_id`` / ``session_id``) is captured by closure via the
  ``make_*`` factories — the same pattern used by the artifacts, Word document,
  and Excel spreadsheet tools (the Strands runtime here does NOT populate
  ``ToolContext.invocation_state`` with identity). The tools are injected
  per-request through ``extra_tools`` (see ``_build_powerpoint_presentation_tools``
  in ``apis/inference_api/chat/routes.py``); they are deliberately NOT registered
  in ``builtin_tools/__init__`` because they need request-scoped identity.
"""

from __future__ import annotations

import asyncio
import logging
from typing import Any, Dict, Optional

from strands import tool

from agents.builtin_tools.office._storage import (
    _DocGenError,
    _ci_exec,
    _ci_read_bytes,
    _ci_write_bytes,
    _download_card,
    _download_s3_bytes,
    _error,
    _get_code_interpreter_id,
    _NO_CI_MESSAGE,
    _region,
    _storage_configured,
    _store_document,
    _validate_document_name,
)

logger = logging.getLogger(__name__)

# PowerPoint presentation MIME type (matches apis.shared.files.ALLOWED_MIME_TYPES).
_PPTX_MIME = (
    "application/vnd.openxmlformats-officedocument.presentationml.presentation"
)

# Sandbox path used to stage a source presentation loaded from S3.
_SANDBOX_SOURCE = "_source.pptx"

# Sandbox path used to stage a template presentation loaded from S3.
_SANDBOX_TEMPLATE = "_template.pptx"

_NO_STORAGE_MESSAGE = (
    "❌ PowerPoint presentation storage is not configured "
    "(S3_USER_FILES_BUCKET_NAME is not set on the runtime)."
)


# ---------------------------------------------------------------------------
# python-pptx presentation builders (run in Code Interpreter)
# ---------------------------------------------------------------------------


_PPTX_PREAMBLE = (
    "from pptx import Presentation\n"
    "from pptx.util import Inches, Pt, Emu\n"
    "from pptx.dml.color import RGBColor\n"
    "from pptx.enum.text import PP_ALIGN, MSO_ANCHOR\n"
    "from pptx.enum.shapes import MSO_SHAPE\n"
)


def _generate_pptx_bytes(
    code_interpreter_id: str,
    python_code: str,
    filename: str,
    template_bytes: Optional[bytes] = None,
) -> bytes:
    """Build a new .pptx from user code and return its bytes.

    When ``template_bytes`` is given, the deck is built on top of that template
    so generated slides inherit its slide masters, layouts, theme colors, fonts,
    and any master/layout branding (logo, footer). Otherwise a blank 16:9 deck
    is used.

    Blocking (boto3 / Code Interpreter) — call via ``asyncio.to_thread``.
    """
    from bedrock_agentcore.tools.code_interpreter_client import CodeInterpreter

    code_interpreter = CodeInterpreter(_region())
    code_interpreter.start(identifier=code_interpreter_id)
    try:
        if template_bytes is not None:
            # Start from the uploaded template, then strip its own example
            # slides — removing the slide-id references leaves the slide
            # masters and layouts (and their theme/branding) intact, so the
            # deck starts themed-but-empty and the model's code adds fresh
            # slides via the template's layouts (``prs.slide_layouts[...]``).
            # Slide size is inherited from the template (not forced to 16:9).
            _ci_write_bytes(code_interpreter, _SANDBOX_TEMPLATE, template_bytes)
            init = (
                f"prs = Presentation({_SANDBOX_TEMPLATE!r})\n"
                "for _sid in list(prs.slides._sldIdLst):\n"
                "    prs.slides._sldIdLst.remove(_sid)\n"
            )
        else:
            init = (
                "prs = Presentation()\n"
                "prs.slide_width = Inches(13.333)\n"
                "prs.slide_height = Inches(7.5)\n"
            )
        # The user's code operates on a pre-initialized ``prs`` and must not
        # call Presentation()/prs.save() itself — we own the lifecycle.
        _ci_exec(
            code_interpreter,
            (
                f"{_PPTX_PREAMBLE}\n"
                f"{init}\n"
                f"{python_code}\n\n"
                f"prs.save({filename!r})\n"
            ),
        )
        data = _ci_read_bytes(code_interpreter, filename)
        if data is None:
            raise _DocGenError(
                f"Presentation '{filename}' was not produced. Make sure your "
                "code adds slides to `prs`."
            )
        return data
    finally:
        try:
            code_interpreter.stop()
        except Exception:  # pragma: no cover - cleanup best-effort
            pass


def _modify_pptx_bytes(
    code_interpreter_id: str,
    source_bytes: bytes,
    python_code: str,
    output_filename: str,
) -> bytes:
    """Load an existing .pptx, apply user edits, return the new bytes.

    Blocking — call via ``asyncio.to_thread``.
    """
    from bedrock_agentcore.tools.code_interpreter_client import CodeInterpreter

    code_interpreter = CodeInterpreter(_region())
    code_interpreter.start(identifier=code_interpreter_id)
    try:
        _ci_write_bytes(code_interpreter, _SANDBOX_SOURCE, source_bytes)
        _ci_exec(
            code_interpreter,
            (
                f"{_PPTX_PREAMBLE}\n"
                f"prs = Presentation({_SANDBOX_SOURCE!r})\n\n"
                f"{python_code}\n\n"
                f"prs.save({output_filename!r})\n"
            ),
        )
        data = _ci_read_bytes(code_interpreter, output_filename)
        if data is None:
            raise _DocGenError(
                f"Modified presentation '{output_filename}' was not produced."
            )
        return data
    finally:
        try:
            code_interpreter.stop()
        except Exception:  # pragma: no cover - cleanup best-effort
            pass


def _extract_pptx_text(code_interpreter_id: str, source_bytes: bytes) -> str:
    """Extract readable text (per slide: shapes, tables, notes) from a .pptx.

    Blocking — call via ``asyncio.to_thread``.
    """
    from bedrock_agentcore.tools.code_interpreter_client import CodeInterpreter

    code_interpreter = CodeInterpreter(_region())
    code_interpreter.start(identifier=code_interpreter_id)
    try:
        _ci_write_bytes(code_interpreter, _SANDBOX_SOURCE, source_bytes)
        extraction = (
            "from pptx import Presentation\n"
            f"prs = Presentation({_SANDBOX_SOURCE!r})\n"
            "lines = []\n"
            "for i, slide in enumerate(prs.slides):\n"
            "    lines.append('## Slide %d' % (i + 1))\n"
            "    for shape in slide.shapes:\n"
            "        if shape.has_text_frame:\n"
            "            t = shape.text_frame.text.strip()\n"
            "            if t:\n"
            "                lines.append(t)\n"
            "        if shape.has_table:\n"
            "            for row in shape.table.rows:\n"
            "                cells = [c.text.strip() for c in row.cells]\n"
            "                lines.append(' | '.join(cells))\n"
            "    if slide.has_notes_slide:\n"
            "        notes = slide.notes_slide.notes_text_frame.text.strip()\n"
            "        if notes:\n"
            "            lines.append('[Notes] ' + notes)\n"
            "    lines.append('')\n"
            "print('\\n'.join(lines))\n"
        )
        return _ci_exec(code_interpreter, extraction).strip()
    finally:
        try:
            code_interpreter.stop()
        except Exception:  # pragma: no cover - cleanup best-effort
            pass


def _extract_pptx_layouts(code_interpreter_id: str, source_bytes: bytes) -> str:
    """List a presentation's slide layouts (index, name, placeholders).

    Useful before building on a template: the model can see which layout
    indices exist and which placeholders each one exposes, so
    ``prs.slides.add_slide(prs.slide_layouts[i])`` targets the right layout and
    ``slide.placeholders[idx]`` fills the right slot. Blocking — call via
    ``asyncio.to_thread``.
    """
    from bedrock_agentcore.tools.code_interpreter_client import CodeInterpreter

    code_interpreter = CodeInterpreter(_region())
    code_interpreter.start(identifier=code_interpreter_id)
    try:
        _ci_write_bytes(code_interpreter, _SANDBOX_SOURCE, source_bytes)
        extraction = (
            "from pptx import Presentation\n"
            f"prs = Presentation({_SANDBOX_SOURCE!r})\n"
            "layouts = prs.slide_layouts\n"
            "lines = ['Total layouts: %d' % len(layouts)]\n"
            "for i, layout in enumerate(layouts):\n"
            "    phs = []\n"
            "    for ph in layout.placeholders:\n"
            "        phs.append('%d=%s' % (ph.placeholder_format.idx, ph.name))\n"
            "    detail = ', '.join(phs) if phs else '(no placeholders)'\n"
            "    lines.append('[%d] %s | placeholders: %s' % (i, layout.name, detail))\n"
            "print('\\n'.join(lines))\n"
        )
        return _ci_exec(code_interpreter, extraction).strip()
    finally:
        try:
            code_interpreter.stop()
        except Exception:  # pragma: no cover - cleanup best-effort
            pass


# ---------------------------------------------------------------------------
# User-files lookup
# ---------------------------------------------------------------------------


async def _find_powerpoint_presentation(
    user_id: str, session_id: str, presentation_name: str
):
    """Find the newest READY .pptx in this session matching ``presentation_name``.

    Returns the ``FileMetadata`` or ``None``. ``list_session_files`` returns
    newest-first, so the first match is the latest version.
    """
    from apis.shared.files import FileStatus, get_file_upload_repository

    target = (
        presentation_name
        if presentation_name.lower().endswith(".pptx")
        else f"{presentation_name}.pptx"
    )
    files = await get_file_upload_repository().list_session_files(
        session_id, status=FileStatus.READY
    )
    for meta in files:
        if (
            meta.user_id == user_id
            and meta.mime_type == _PPTX_MIME
            and meta.filename.lower() == target.lower()
        ):
            return meta
    return None


# ---------------------------------------------------------------------------
# Tool factories
# ---------------------------------------------------------------------------


def make_create_powerpoint_presentation_tool(session_id: str, user_id: str):
    """Create a ``create_powerpoint_presentation`` tool bound to the identity."""

    @tool
    async def create_powerpoint_presentation(
        python_code: str,
        presentation_name: str,
        template_name: Optional[str] = None,
    ) -> Any:
        """Create a new PowerPoint (.pptx) presentation using python-pptx code.

        Executes python-pptx code in a sandboxed Code Interpreter to build a
        16:9 widescreen deck, saves it to the user's files, and returns a
        download card. Great for pitch decks, reports, and summaries with
        titled slides, bullet content, tables, and embedded charts.

        Optionally builds on a template (see ``template_name``) so the deck
        inherits a branded theme, fonts, and layouts instead of the plain
        default. Prefer a template when the user has uploaded one or asked for
        their branding.

        Available libraries in the sandbox: python-pptx, matplotlib, pandas,
        numpy.

        Args:
            python_code: python-pptx code that builds the deck. A blank 16:9
                presentation is already available as ``prs = Presentation()``
                (slide size preset to 13.33" x 7.5") — do NOT call
                ``Presentation()`` or ``prs.save()`` yourself; the tool saves it
                for you. ``Inches``, ``Pt``, ``Emu``, ``RGBColor``, ``PP_ALIGN``,
                ``MSO_ANCHOR`` and ``MSO_SHAPE`` are already imported.

                Add slides from the built-in layouts (0=title, 1=title+content,
                5=title only, 6=blank), e.g.:
                    slide = prs.slides.add_slide(prs.slide_layouts[0])
                    slide.shapes.title.text = 'Quarterly Review'
                    slide.placeholders[1].text = 'FY2026 — Q4'

                A dark title slide with a custom text box:
                    slide = prs.slides.add_slide(prs.slide_layouts[6])
                    slide.background.fill.solid()
                    slide.background.fill.fore_color.rgb = RGBColor(0x1E, 0x27, 0x61)
                    box = slide.shapes.add_textbox(Inches(0.8), Inches(2.6), Inches(11.7), Inches(2))
                    tf = box.text_frame
                    tf.text = 'Product Strategy'
                    r = tf.paragraphs[0].runs[0]
                    r.font.size, r.font.bold = Pt(44), True
                    r.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

                A bullet content slide (keep to <= 4 bullets):
                    slide = prs.slides.add_slide(prs.slide_layouts[1])
                    slide.shapes.title.text = 'Highlights'
                    body = slide.placeholders[1].text_frame
                    body.text = 'Revenue up 15%'
                    p = body.add_paragraph(); p.text = 'Churn down to 2%'

                A table:
                    tbl = slide.shapes.add_table(2, 2, Inches(1), Inches(2),
                                                 Inches(8), Inches(2)).table
                    tbl.cell(0, 0).text = 'Quarter'; tbl.cell(0, 1).text = 'Revenue'

                A matplotlib chart image:
                    import matplotlib.pyplot as plt
                    plt.figure(figsize=(8, 4.5))
                    plt.bar(['Q1', 'Q2'], [100, 120])
                    plt.savefig('chart.png', dpi=200, bbox_inches='tight')
                    plt.close()
                    slide.shapes.add_picture('chart.png', Inches(2.5), Inches(1.5), width=Inches(8))

                Design guidance (aim for a polished, cohesive deck):
                - Pick ONE color palette for the whole deck and reuse it. Good
                  options (primary / secondary / accent hex):
                  Midnight Executive 1E2761 / CADCFC / FFFFFF;
                  Teal Trust 028090 / 00A896 / 02C39A;
                  Forest & Moss 2C5F2D / 97BC62 / F5F5F5;
                  Coral Energy F96167 / F9E795 / 2F3C7E;
                  Charcoal Minimal 36454F / F2F2F2 / 212121.
                  One color dominates (60-70%); dark backgrounds for the
                  title/closing slides, light for content. Never plain white,
                  never default PowerPoint blue.
                - Typography: titles 36-44pt bold, section headers 20-24pt,
                  body 14-16pt, big stat callouts 48-120pt. Left-align body
                  text; center only titles and stats.
                - Every slide should carry a visual element (a shape, colored
                  accent bar, icon circle, table, or chart) — avoid text-only
                  slides, and vary the layout slide to slide.

            presentation_name: File name WITHOUT extension (.pptx is added
                automatically). Use only letters, numbers, hyphens, and
                underscores (e.g. "sales-deck", "Q4_review").

            template_name: Optional name of a .pptx template already available
                in this chat (an uploaded deck or a previously generated one,
                with or without the .pptx extension). When given, the new deck
                is built ON TOP of that template so it inherits the template's
                slide masters, layouts, theme colors, fonts, and any master- or
                layout-level branding (logo, footer). The template's own example
                slides are stripped first, so you start themed-but-empty.

                When using a template:
                - Add slides from the TEMPLATE's layouts, e.g.
                  ``slide = prs.slides.add_slide(prs.slide_layouts[1])``, and
                  fill the layout's placeholders (``slide.shapes.title``,
                  ``slide.placeholders[idx]``) rather than drawing everything
                  from scratch — that's what preserves the branded look.
                - Do NOT override slide backgrounds/fonts with the palette
                  below; let the template's theme drive the design.
                - Use ``list_powerpoint_presentations`` to see available names,
                  and ``read_powerpoint_presentation`` to inspect a template's
                  existing content if helpful.

        Returns:
            An inline download card. The presentation is also saved to this
            chat's Files.
        """
        is_valid, error_msg = _validate_document_name(presentation_name)
        if not is_valid:
            return _error(
                f"❌ Invalid presentation name '{presentation_name}': {error_msg}\n\n"
                "Examples: sales-deck, Q4_review, pitch-final"
            )

        filename = f"{presentation_name}.pptx"
        code_interpreter_id = _get_code_interpreter_id()
        if not code_interpreter_id:
            return _error(_NO_CI_MESSAGE)
        if not _storage_configured():
            return _error(_NO_STORAGE_MESSAGE)

        # Resolve an optional template from this chat's files. When provided,
        # the deck is built on top of it so it inherits the template's theme.
        template_bytes = None
        if template_name:
            template_src = await _find_powerpoint_presentation(
                user_id, session_id, template_name
            )
            if template_src is None:
                return _error(
                    f"❌ No PowerPoint template named '{template_name}' was found "
                    "in this chat. Upload a .pptx template first, or use "
                    "list_powerpoint_presentations to see what's available."
                )
            try:
                template_bytes = await asyncio.to_thread(
                    _download_s3_bytes, template_src.s3_bucket, template_src.s3_key
                )
            except Exception as exc:  # noqa: BLE001 - surface storage errors
                logger.error(f"create_powerpoint_presentation template load error: {exc}")
                return _error(
                    f"❌ Failed to load template '{template_src.filename}': {exc}"
                )

        try:
            file_bytes = await asyncio.to_thread(
                _generate_pptx_bytes,
                code_interpreter_id,
                python_code,
                filename,
                template_bytes,
            )
        except _DocGenError as exc:
            return _error(
                f"❌ Failed to create '{filename}'.\n\n```\n{exc}\n```\n\n"
                "Check the python-pptx code for errors."
            )
        except Exception as exc:  # noqa: BLE001 - surface any sandbox error
            logger.error(f"create_powerpoint_presentation sandbox error: {exc}")
            return _error(f"❌ Failed to create '{filename}': {exc}")

        try:
            _id, download_url, size_kb = await _store_document(
                user_id, session_id, filename, file_bytes, _PPTX_MIME
            )
        except Exception as exc:  # noqa: BLE001 - storage failure is terminal
            logger.error(f"create_powerpoint_presentation storage error: {exc}")
            return _error(f"❌ Created '{filename}' but failed to save it: {exc}")

        return _download_card(filename, download_url, size_kb, "Created")

    return create_powerpoint_presentation


def make_modify_powerpoint_presentation_tool(session_id: str, user_id: str):
    """Create a ``modify_powerpoint_presentation`` tool bound to the identity."""

    @tool
    async def modify_powerpoint_presentation(
        presentation_name: str,
        python_code: str,
        output_name: Optional[str] = None,
    ) -> Any:
        """Modify an existing PowerPoint (.pptx) presentation with python-pptx code.

        Loads a deck previously created in this chat, runs your python-pptx code
        against it, and saves the result (as a new file so the original is
        preserved). Returns a download card.

        Use ``list_powerpoint_presentations`` first if you are unsure of the
        exact name.

        Args:
            presentation_name: Name of the existing deck to edit (with or
                without the .pptx extension), e.g. "sales-deck".
            python_code: python-pptx code that edits the deck. The loaded
                presentation is available as ``prs = Presentation(...)`` — do
                NOT call ``Presentation()`` or ``prs.save()`` yourself. Existing
                slides are ``prs.slides``; add new ones with
                ``prs.slides.add_slide(prs.slide_layouts[...])``. ``Inches``,
                ``Pt``, ``Emu``, ``RGBColor``, ``PP_ALIGN``, ``MSO_ANCHOR`` and
                ``MSO_SHAPE`` are already imported.

                Example (append a closing slide):
                    slide = prs.slides.add_slide(prs.slide_layouts[5])
                    slide.shapes.title.text = 'Thank You'

                Example (edit the first slide's title):
                    prs.slides[0].shapes.title.text = 'Updated Title'

            output_name: Optional name (without extension) for the edited copy.
                Defaults to the source name (a new versioned copy is saved).

        Returns:
            An inline download card for the edited presentation.
        """
        code_interpreter_id = _get_code_interpreter_id()
        if not code_interpreter_id:
            return _error(_NO_CI_MESSAGE)
        if not _storage_configured():
            return _error(_NO_STORAGE_MESSAGE)

        source = await _find_powerpoint_presentation(
            user_id, session_id, presentation_name
        )
        if source is None:
            return _error(
                f"❌ No PowerPoint presentation named '{presentation_name}' was "
                "found in this chat. Use list_powerpoint_presentations to see "
                "what's available."
            )

        out_base = output_name or source.filename
        if out_base.lower().endswith(".pptx"):
            out_base = out_base[: -len(".pptx")]
        is_valid, error_msg = _validate_document_name(out_base)
        if not is_valid:
            return _error(
                f"❌ Invalid output name '{out_base}': {error_msg}"
            )
        output_filename = f"{out_base}.pptx"

        try:
            source_bytes = await asyncio.to_thread(
                _download_s3_bytes, source.s3_bucket, source.s3_key
            )
            file_bytes = await asyncio.to_thread(
                _modify_pptx_bytes,
                code_interpreter_id,
                source_bytes,
                python_code,
                output_filename,
            )
        except _DocGenError as exc:
            return _error(
                f"❌ Failed to modify '{source.filename}'.\n\n```\n{exc}\n```\n\n"
                "Check the python-pptx code for errors."
            )
        except Exception as exc:  # noqa: BLE001 - surface any sandbox error
            logger.error(f"modify_powerpoint_presentation error: {exc}")
            return _error(f"❌ Failed to modify '{source.filename}': {exc}")

        try:
            _id, download_url, size_kb = await _store_document(
                user_id, session_id, output_filename, file_bytes, _PPTX_MIME
            )
        except Exception as exc:  # noqa: BLE001 - storage failure is terminal
            logger.error(f"modify_powerpoint_presentation storage error: {exc}")
            return _error(
                f"❌ Modified '{source.filename}' but failed to save it: {exc}"
            )

        return _download_card(output_filename, download_url, size_kb, "Updated")

    return modify_powerpoint_presentation


def make_list_powerpoint_presentations_tool(session_id: str, user_id: str):
    """Create a ``list_powerpoint_presentations`` tool bound to the identity."""

    @tool
    async def list_powerpoint_presentations() -> Dict[str, Any]:
        """List the PowerPoint (.pptx) presentations available in this chat.

        Returns the file names and sizes of decks created or modified in this
        conversation. Use the names with modify_powerpoint_presentation or
        read_powerpoint_presentation.
        """
        from apis.shared.files import FileStatus, get_file_upload_repository

        files = await get_file_upload_repository().list_session_files(
            session_id, status=FileStatus.READY
        )
        seen: set[str] = set()
        rows = []
        for meta in files:  # newest-first
            if meta.user_id != user_id or meta.mime_type != _PPTX_MIME:
                continue
            if meta.filename in seen:
                continue
            seen.add(meta.filename)
            rows.append(f"- {meta.filename} ({meta.size_bytes / 1024:.1f} KB)")

        if not rows:
            text = (
                "No PowerPoint presentations in this chat yet. Use "
                "create_powerpoint_presentation to make one."
            )
        else:
            text = "PowerPoint presentations in this chat:\n" + "\n".join(rows)
        return {"content": [{"text": text}], "status": "success"}

    return list_powerpoint_presentations


def make_read_powerpoint_presentation_tool(session_id: str, user_id: str):
    """Create a ``read_powerpoint_presentation`` tool bound to the identity."""

    @tool
    async def read_powerpoint_presentation(presentation_name: str) -> Dict[str, Any]:
        """Read the text content of an existing PowerPoint (.pptx) presentation.

        Extracts each slide's text, tables, and speaker notes from a deck
        created in this chat so you can reference or summarize its contents. Use
        list_powerpoint_presentations first if unsure of the exact name.

        Args:
            presentation_name: Name of the deck to read (with or without the
                .pptx extension), e.g. "sales-deck".

        Returns:
            The presentation's text content, grouped by slide.
        """
        code_interpreter_id = _get_code_interpreter_id()
        if not code_interpreter_id:
            return _error(_NO_CI_MESSAGE)
        if not _storage_configured():
            return _error(_NO_STORAGE_MESSAGE)

        source = await _find_powerpoint_presentation(
            user_id, session_id, presentation_name
        )
        if source is None:
            return _error(
                f"❌ No PowerPoint presentation named '{presentation_name}' was "
                "found in this chat. Use list_powerpoint_presentations to see "
                "what's available."
            )

        try:
            source_bytes = await asyncio.to_thread(
                _download_s3_bytes, source.s3_bucket, source.s3_key
            )
            text = await asyncio.to_thread(
                _extract_pptx_text, code_interpreter_id, source_bytes
            )
        except _DocGenError as exc:
            return _error(f"❌ Failed to read '{source.filename}': {exc}")
        except Exception as exc:  # noqa: BLE001 - surface any sandbox error
            logger.error(f"read_powerpoint_presentation error: {exc}")
            return _error(f"❌ Failed to read '{source.filename}': {exc}")

        body = text or "(The presentation has no extractable text.)"
        return {
            "content": [
                {"text": f"Content of {source.filename}:\n\n{body}"}
            ],
            "status": "success",
        }

    return read_powerpoint_presentation


def make_list_powerpoint_layouts_tool(session_id: str, user_id: str):
    """Create a ``list_powerpoint_layouts`` tool bound to the identity."""

    @tool
    async def list_powerpoint_layouts(presentation_name: str) -> Dict[str, Any]:
        """List the slide layouts of a PowerPoint (.pptx) file in this chat.

        Reports each layout's index, name, and placeholder slots. Call this on a
        template BEFORE building on it (create_powerpoint_presentation with
        template_name) so you add slides from the right layout
        (``prs.slide_layouts[index]``) and fill the correct placeholders
        (``slide.placeholders[idx]``) — that's what preserves the template's
        branded design. Works on any .pptx available in this chat (an uploaded
        template or a deck generated here).

        Args:
            presentation_name: Name of the .pptx to inspect (with or without the
                .pptx extension), e.g. "brand-template".

        Returns:
            The layout inventory (index, name, placeholders per layout).
        """
        code_interpreter_id = _get_code_interpreter_id()
        if not code_interpreter_id:
            return _error(_NO_CI_MESSAGE)
        if not _storage_configured():
            return _error(_NO_STORAGE_MESSAGE)

        source = await _find_powerpoint_presentation(
            user_id, session_id, presentation_name
        )
        if source is None:
            return _error(
                f"❌ No PowerPoint presentation named '{presentation_name}' was "
                "found in this chat. Use list_powerpoint_presentations to see "
                "what's available."
            )

        try:
            source_bytes = await asyncio.to_thread(
                _download_s3_bytes, source.s3_bucket, source.s3_key
            )
            text = await asyncio.to_thread(
                _extract_pptx_layouts, code_interpreter_id, source_bytes
            )
        except _DocGenError as exc:
            return _error(f"❌ Failed to inspect '{source.filename}': {exc}")
        except Exception as exc:  # noqa: BLE001 - surface any sandbox error
            logger.error(f"list_powerpoint_layouts error: {exc}")
            return _error(f"❌ Failed to inspect '{source.filename}': {exc}")

        body = text or "(No layouts found.)"
        return {
            "content": [
                {"text": f"Layouts in {source.filename}:\n\n{body}"}
            ],
            "status": "success",
        }

    return list_powerpoint_layouts
